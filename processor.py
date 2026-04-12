#!/usr/bin/env python3
"""
Pitch Deck Email Processor
Checks Gmail daily for emails containing pitch decks (funding round decks),
analyses them with Claude, saves to Google Drive, and updates Attio CRM.

SETUP: See README.md for first-time setup instructions.
"""

import os
import sys
import io
import re
import json
import base64
import pickle
import email.mime.text
from datetime import datetime, timedelta, timezone
from pathlib import Path

import requests as _requests

# Gmail & Drive API
from google.auth.transport.requests import Request
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseUpload

# Claude
import anthropic

# Attio CRM
import attio

# File extraction
try:
    from pypdf import PdfReader
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False

try:
    from pptx import Presentation
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False

# ── Config ─────────────────────────────────────────────────────────────────────

SCRIPT_DIR = Path(__file__).parent
CREDENTIALS_FILE = SCRIPT_DIR / "credentials.json"
TOKEN_FILE = SCRIPT_DIR / "token.pickle"
PROCESSED_LOG = SCRIPT_DIR / "processed.json"

# Google Drive folder IDs where pitch decks will be saved simultaneously
DRIVE_FOLDER_IDS = [
    "1bg0NQVwuP82wkIWvXzlJCrs-WHYk12DD",  # ar@angelinvest.ventures
    "13CKApFKyLmlcl90Sa-xEXMcaqHnkz3tB",  # anna.ritz@legata.cc
]

NOTIFICATION_EMAIL = "ar@angelinvest.ventures"

# Email to pass to DocSend for access-controlled decks (the address decks are shared with)
DOCSEND_EMAIL = "ar@angelinvest.ventures"

# Regex to find DocSend links in email bodies
_DOCSEND_RE = re.compile(r'https?://(?:www\.)?docsend\.com/view/[a-zA-Z0-9]+(?:/[a-zA-Z0-9]+)?')

# Only attachments with these MIME types are considered
PITCH_DECK_MIME_TYPES = {
    "application/pdf",
    "application/vnd.ms-powerpoint",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

SCOPES = [
    "https://www.googleapis.com/auth/gmail.readonly",
    "https://www.googleapis.com/auth/gmail.send",
    "https://www.googleapis.com/auth/drive",
]

# ── Auth ───────────────────────────────────────────────────────────────────────

def get_credentials():
    """Load or refresh OAuth credentials, prompting via console if needed."""
    creds = None
    if TOKEN_FILE.exists():
        with open(TOKEN_FILE, "rb") as f:
            creds = pickle.load(f)

    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            if not CREDENTIALS_FILE.exists():
                print("ERROR: credentials.json not found.")
                print("Follow the setup guide: see README.md in this folder.")
                sys.exit(1)
            flow = InstalledAppFlow.from_client_secrets_file(
                str(CREDENTIALS_FILE), SCOPES
            )
            # Force correct redirect_uri to avoid credentials.json copy-paste corruption
            flow.redirect_uri = "http://localhost"
            # Console-based flow — no localhost server required (works from any device)
            auth_url, _ = flow.authorization_url(prompt="consent")
            print("\nAuthorisation required. Visit this URL in your browser:")
            print(f"\n  {auth_url}\n")
            code = input("Paste the authorisation code here: ").strip()
            flow.fetch_token(code=code)
            creds = flow.credentials

        with open(TOKEN_FILE, "wb") as f:
            pickle.dump(creds, f)

    return creds


def get_gmail_service():
    return build("gmail", "v1", credentials=get_credentials())


def get_drive_service():
    return build("drive", "v3", credentials=get_credentials())


# ── Drive helpers ──────────────────────────────────────────────────────────────

def upload_to_drive(drive_service, folder_ids: list, filename: str, data: bytes, mime_type: str) -> dict:
    """Upload bytes to multiple Drive folders simultaneously. Returns the file resource."""
    meta = {"name": filename, "parents": folder_ids}
    media = MediaIoBaseUpload(io.BytesIO(data), mimetype=mime_type, resumable=False)
    return drive_service.files().create(
        body=meta, media_body=media, fields="id, name, webViewLink"
    ).execute()


# ── Processed log ──────────────────────────────────────────────────────────────

def load_processed() -> set:
    if PROCESSED_LOG.exists():
        with open(PROCESSED_LOG) as f:
            return set(json.load(f))
    return set()


def save_processed(processed: set):
    with open(PROCESSED_LOG, "w") as f:
        json.dump(list(processed), f, indent=2)


# ── Text extraction ────────────────────────────────────────────────────────────

def extract_text_from_pdf(data: bytes) -> str:
    if not PDF_SUPPORT:
        return "[PDF text extraction unavailable — install pypdf]"
    try:
        reader = PdfReader(io.BytesIO(data))
        pages = [page.extract_text() or "" for page in reader.pages[:20]]
        return "\n".join(pages)[:8000]
    except Exception as e:
        return f"[PDF extraction error: {e}]"


def extract_text_from_pptx(data: bytes) -> str:
    if not PPTX_SUPPORT:
        return "[PPTX text extraction unavailable — install python-pptx]"
    try:
        prs = Presentation(io.BytesIO(data))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)[:8000]
    except Exception as e:
        return f"[PPTX extraction error: {e}]"


def extract_text(data: bytes, mime_type: str) -> str:
    if mime_type == "application/pdf":
        return extract_text_from_pdf(data)
    if mime_type in (
        "application/vnd.ms-powerpoint",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ):
        return extract_text_from_pptx(data)
    return ""


# ── Claude analysis ────────────────────────────────────────────────────────────

def analyze_deck(email_subject: str, email_body: str, attachment_name: str, attachment_text: str) -> dict:
    """
    Ask Claude whether this is a pitch deck and extract company information.
    Returns a dict with: is_pitch_deck, confidence, reasoning,
                         company_name, trade_name, domain, founders
    """
    client = anthropic.Anthropic()

    prompt = f"""You are analysing an email to determine whether it contains a startup or company funding pitch deck (i.e. a deck used to raise investment). If it is a pitch deck, also extract key company information.

EMAIL SUBJECT:
{email_subject}

EMAIL BODY (first 2000 chars):
{email_body[:2000]}

ATTACHMENT FILENAME:
{attachment_name}

ATTACHMENT TEXT (first 4000 chars):
{attachment_text[:4000]}

Respond with a JSON object only, no other text:
{{
  "is_pitch_deck": true or false,
  "confidence": "high" | "medium" | "low",
  "reasoning": "one sentence explanation",
  "company_name": "official/legal company name, or null if not a pitch deck",
  "trade_name": "brand/trading name if different from company_name, else null",
  "domain": "company website domain e.g. sprive.com (no https://), or null if not found",
  "founders": ["founder full name 1", "founder full name 2"]
}}"""

    message = client.messages.create(
        model="claude-opus-4-6",
        max_tokens=512,
        messages=[{"role": "user", "content": prompt}],
    )

    text = message.content[0].text.strip()
    if text.startswith("```"):
        text = text.split("```")[1]
        if text.startswith("json"):
            text = text[4:]
    return json.loads(text)


# ── Email notifications ────────────────────────────────────────────────────────

def send_notification_email(gmail_service, subject: str, body: str):
    """Send a notification email via Gmail API."""
    msg = email.mime.text.MIMEText(body)
    msg["to"] = NOTIFICATION_EMAIL
    msg["from"] = NOTIFICATION_EMAIL
    msg["subject"] = subject
    raw = base64.urlsafe_b64encode(msg.as_bytes()).decode()
    gmail_service.users().messages().send(userId="me", body={"raw": raw}).execute()
    print(f"    Notification sent: {subject}")


# ── Gmail helpers ──────────────────────────────────────────────────────────────

def get_email_body(payload) -> str:
    if payload.get("mimeType") == "text/plain":
        data = payload.get("body", {}).get("data", "")
        return base64.urlsafe_b64decode(data + "==").decode("utf-8", errors="replace")
    for part in payload.get("parts", []):
        body = get_email_body(part)
        if body:
            return body
    return ""


def download_attachment(service, message_id: str, attachment_id: str) -> bytes:
    attachment = service.users().messages().attachments().get(
        userId="me", messageId=message_id, id=attachment_id
    ).execute()
    data = attachment.get("data", "")
    return base64.urlsafe_b64decode(data + "==")


def sanitise_filename(name: str) -> str:
    return "".join(c for c in name if c not in r'\/:*?"<>|').strip()


# ── Attio integration ──────────────────────────────────────────────────────────

def handle_attio(gmail_service, analysis: dict, drive_link: str, email_subject: str, sender: str):
    """Match company in Attio and create/update record. Never deletes anything."""
    company_name = analysis.get("company_name") or "Unknown Company"
    trade_name = analysis.get("trade_name")
    domain = analysis.get("domain")
    founders = analysis.get("founders", [])

    print(f"    Attio: matching '{company_name}' (domain: {domain}, trade name: {trade_name})")

    status, candidates = attio.match_company(company_name, domain, trade_name)

    if status == "single_match":
        record = candidates[0]
        record_id = attio.get_record_id(record)
        attio_name = attio.get_company_name(record)

        if attio.is_owned_by_anna_ritz(record):
            attio.update_pitch_deck_url(record_id, drive_link)
            print(f"    Attio: updated '{attio_name}' with pitch deck URL")
        else:
            subject = f"Pitch deck received — {company_name} (owner is not you)"
            body = (
                f"A pitch deck was received and saved to Drive, but the Attio record is owned by someone else.\n\n"
                f"Company: {company_name}\n"
                f"Attio record: {attio_name}\n"
                f"Email subject: {email_subject}\n"
                f"Sender: {sender}\n"
                f"Drive link: {drive_link}\n\n"
                f"Please update Attio manually if needed."
            )
            send_notification_email(gmail_service, subject, body)

    elif status == "no_match":
        attio.create_company(company_name, domain, drive_link)
        print(f"    Attio: created new record for '{company_name}'")

    elif status == "ambiguous":
        candidate_names = [attio.get_company_name(c) for c in candidates]
        subject = f"Pitch deck received — ambiguous Attio match: {company_name}"
        body = (
            f"A pitch deck was received and saved to Drive, but Attio matching was ambiguous.\n\n"
            f"Extracted company name: {company_name}\n"
            f"Trade name: {trade_name or 'N/A'}\n"
            f"Domain: {domain or 'N/A'}\n"
            f"Founders: {', '.join(founders) or 'N/A'}\n"
            f"Email subject: {email_subject}\n"
            f"Sender: {sender}\n"
            f"Drive link: {drive_link}\n\n"
            f"Attio candidates found:\n" +
            "\n".join(f"  - {name}" for name in candidate_names) +
            "\n\nPlease update Attio manually."
        )
        send_notification_email(gmail_service, subject, body)


# ── DocSend helpers ────────────────────────────────────────────────────────────

def extract_docsend_links(text: str) -> list:
    """Return deduplicated list of DocSend URLs found in text."""
    return list(dict.fromkeys(_DOCSEND_RE.findall(text)))


def download_docsend_pdf(docsend_url: str) -> bytes | None:
    """
    Download a DocSend deck as PDF via docsend2pdf.com API.
    Returns raw PDF bytes, or None if the download fails.
    """
    try:
        r = _requests.post(
            "https://docsend2pdf.com/api/convert",
            json={"url": docsend_url, "email": DOCSEND_EMAIL},
            timeout=90,
        )
        if r.status_code == 200:
            return r.content
        print(f"    DocSend API error {r.status_code}: {r.text[:300]}")
        return None
    except Exception as e:
        print(f"    DocSend download error: {e}")
        return None


# ── Main processing ────────────────────────────────────────────────────────────

def process_emails():
    gmail_service = get_gmail_service()
    drive_service = get_drive_service()
    attio.initialise()
    processed = load_processed()

    # Look back 25 hours to cover the daily run with overlap
    cutoff = datetime.now(timezone.utc) - timedelta(hours=25)
    after_epoch = int(cutoff.timestamp())
    query = f"(has:attachment OR docsend.com/view) after:{after_epoch}"

    print(f"[{datetime.now().strftime('%Y-%m-%d %H:%M')}] Searching Gmail: {query}")

    results = gmail_service.users().messages().list(
        userId="me", q=query, maxResults=50
    ).execute()

    messages = results.get("messages", [])
    print(f"  Found {len(messages)} emails with attachments in the last 25 hours.")

    saved_count = 0

    for msg_ref in messages:
        msg_id = msg_ref["id"]

        if msg_id in processed:
            continue

        msg = gmail_service.users().messages().get(
            userId="me", id=msg_id, format="full"
        ).execute()

        headers = {h["name"]: h["value"] for h in msg["payload"].get("headers", [])}
        subject = headers.get("Subject", "(no subject)")
        sender = headers.get("From", "unknown")
        body = get_email_body(msg["payload"])

        print(f"\n  Email: '{subject}' from {sender}")

        parts = msg["payload"].get("parts", [])
        for part in parts:
            filename = part.get("filename", "")
            mime_type = part.get("mimeType", "")
            attachment_id = part.get("body", {}).get("attachmentId")

            if not attachment_id or mime_type not in PITCH_DECK_MIME_TYPES:
                continue

            print(f"    Attachment: {filename} ({mime_type})")

            attachment_data = download_attachment(gmail_service, msg_id, attachment_id)
            attachment_text = extract_text(attachment_data, mime_type)

            analysis = analyze_deck(subject, body, filename, attachment_text)
            print(f"    Claude says pitch deck: {analysis.get('is_pitch_deck')} — {analysis.get('reasoning')}")

            if analysis.get("is_pitch_deck"):
                safe_name = sanitise_filename(filename)
                timestamp = datetime.now().strftime("%Y%m%d_%H%M")
                upload_name = f"{timestamp}_{safe_name}"

                uploaded = upload_to_drive(
                    drive_service, DRIVE_FOLDER_IDS, upload_name, attachment_data, mime_type
                )
                drive_link = uploaded.get("webViewLink", "")
                print(f"    Saved to Drive: {uploaded.get('name')} — {drive_link}")
                saved_count += 1

                if os.environ.get("ATTIO_API_KEY"):
                    try:
                        handle_attio(gmail_service, analysis, drive_link, subject, sender)
                    except Exception as e:
                        print(f"    Attio error (non-fatal): {e}")

        # ── DocSend links in email body ──────────────────────────────────
        docsend_links = extract_docsend_links(body)
        for ds_url in docsend_links:
            print(f"    DocSend link found: {ds_url}")
            pdf_data = download_docsend_pdf(ds_url)
            if not pdf_data:
                send_notification_email(
                    gmail_service,
                    f"DocSend download failed — {subject}",
                    (
                        f"A DocSend link was found in an email but could not be downloaded.\n\n"
                        f"Link: {ds_url}\n"
                        f"Email subject: {subject}\n"
                        f"Sender: {sender}\n\n"
                        f"The deck may be password-protected or require a specific email address."
                    ),
                )
                continue

            # Use the URL slug as the filename
            slug = ds_url.rstrip("/").split("/")[-1]
            filename = f"docsend_{slug}.pdf"
            attachment_text = extract_text_from_pdf(pdf_data)

            analysis = analyze_deck(subject, body, filename, attachment_text)
            print(f"    Claude says pitch deck: {analysis.get('is_pitch_deck')} — {analysis.get('reasoning')}")

            if analysis.get("is_pitch_deck"):
                timestamp = datetime.now().strftime("%Y%m%d_%H%M")
                upload_name = f"{timestamp}_{filename}"

                uploaded = upload_to_drive(
                    drive_service, DRIVE_FOLDER_IDS, upload_name, pdf_data, "application/pdf"
                )
                drive_link = uploaded.get("webViewLink", "")
                print(f"    Saved to Drive: {uploaded.get('name')} — {drive_link}")
                saved_count += 1

                if os.environ.get("ATTIO_API_KEY"):
                    try:
                        handle_attio(gmail_service, analysis, drive_link, subject, sender)
                    except Exception as e:
                        print(f"    Attio error (non-fatal): {e}")

        processed.add(msg_id)

    save_processed(processed)
    print(f"\nDone. Saved {saved_count} pitch deck(s).")


if __name__ == "__main__":
    process_emails()
